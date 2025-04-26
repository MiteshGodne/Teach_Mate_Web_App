# app/utils/ppt_parser.py
from pptx import Presentation

def extract_ppt_content(file_path):
    prs = Presentation(file_path)
    slides_content = []

    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        slides_content.append({"text": text.strip()})
    return slides_content
